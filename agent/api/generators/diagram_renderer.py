"""
AGRA — Diagram Renderer (Phase 2)
Renders structured diagram JSON as native PowerPoint shapes.

Supports:
  - flowchart:     Linear/branching process flows
  - hierarchy:     Org chart / tree (top-down)
  - block_diagram: System architecture blocks (left-to-right)
  - cycle:         Circular process diagrams

All rendering uses python-pptx AutoShapes — zero external dependencies.
"""

import logging
import math
from typing import Any, Dict, List, Optional, Tuple

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("agra.diagram_renderer")

# ── ICG Palette for diagrams ──
_NAVY = RGBColor(0x0B, 0x10, 0x20)
_DARK_BLUE = RGBColor(0x1E, 0x3A, 0x8A)
_MID_BLUE = RGBColor(0x2C, 0x5F, 0xC9)
_GOLD = RGBColor(0xD4, 0xA5, 0x37)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
_ACCENT_TEAL = RGBColor(0x06, 0xB6, 0xD4)
_ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
_ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
_ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)

_NODE_COLORS = [_DARK_BLUE, _MID_BLUE, _ACCENT_TEAL, _ACCENT_GREEN, _ACCENT_AMBER, _ACCENT_RED]

# Slide working area (inside margins)
_LEFT_MARGIN = 0.5   # inches
_TOP_MARGIN = 1.6    # below title bar
_RIGHT_MARGIN = 0.5
_BOTTOM_MARGIN = 0.5
_SLIDE_W = 10.0
_SLIDE_H = 7.5
_WORK_W = _SLIDE_W - _LEFT_MARGIN - _RIGHT_MARGIN  # 9.0"
_WORK_H = _SLIDE_H - _TOP_MARGIN - _BOTTOM_MARGIN   # 5.4"


def _shape_type_for(shape_name: str) -> int:
    """Map diagram node shape name to MSO_SHAPE enum."""
    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "oval": MSO_SHAPE.OVAL,
        "hexagon": MSO_SHAPE.HEXAGON,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "cloud": MSO_SHAPE.CLOUD,
        "cylinder": MSO_SHAPE.CAN,
        "start_end": MSO_SHAPE.ROUNDED_RECTANGLE,
        "process": MSO_SHAPE.RECTANGLE,
        "decision": MSO_SHAPE.DIAMOND,
    }
    return mapping.get(shape_name, MSO_SHAPE.ROUNDED_RECTANGLE)


def _add_shape_with_text(
    slide,
    shape_type: int,
    left: float, top: float,
    width: float, height: float,
    text: str,
    fill_color: RGBColor = _DARK_BLUE,
    text_color: RGBColor = _WHITE,
    font_size: int = 10,
    bold: bool = False,
):
    """Add a positioned shape with centered text."""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = _GOLD
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Vertical centering
    try:
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER

    return shape


def _add_connector(slide, x1: float, y1: float, x2: float, y2: float, label: str = ""):
    """Add a line connector between two points (in inches)."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = _GOLD
    connector.line.width = Pt(2)

    # Add arrowhead
    try:
        connector.line.end_marker_style = 1  # Triangle arrowhead
    except Exception:
        pass

    # Add label if provided
    if label:
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        txBox = slide.shapes.add_textbox(
            Inches(mid_x - 0.3), Inches(mid_y - 0.15),
            Inches(0.6), Inches(0.3),
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.color.rgb = _LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER


def _layout_flowchart(nodes: list, edges: list) -> Dict[str, Tuple[float, float]]:
    """Auto-layout for flowchart: arrange nodes in a top-to-bottom grid."""
    n = len(nodes)
    if n == 0:
        return {}

    # Determine grid: prefer vertical flow
    cols = min(3, max(1, n // 4 + 1))
    rows = math.ceil(n / cols)

    node_w = min(2.2, _WORK_W / cols - 0.3)
    node_h = 0.8
    h_gap = (_WORK_W - cols * node_w) / max(cols + 1, 1)
    v_gap = min(1.0, (_WORK_H - rows * node_h) / max(rows + 1, 1))

    positions = {}
    for i, node in enumerate(nodes):
        row = i // cols
        col = i % cols
        x = _LEFT_MARGIN + h_gap * (col + 1) + node_w * col
        y = _TOP_MARGIN + v_gap * (row + 1) + node_h * row
        positions[node.get("id", str(i))] = (x, y)

    return positions, node_w, node_h


def _layout_hierarchy(nodes: list, edges: list) -> Dict[str, Tuple[float, float]]:
    """Auto-layout for hierarchy: root at top, children below."""
    # Build parent→children map
    children_map = {}
    child_set = set()
    for e in edges:
        parent = e.get("from", "")
        child = e.get("to", "")
        children_map.setdefault(parent, []).append(child)
        child_set.add(child)

    # Find root(s)
    node_ids = [n.get("id", str(i)) for i, n in enumerate(nodes)]
    roots = [nid for nid in node_ids if nid not in child_set]
    if not roots:
        roots = [node_ids[0]] if node_ids else []

    # BFS to assign levels
    levels = {}
    queue = [(r, 0) for r in roots]
    visited = set()
    while queue:
        nid, level = queue.pop(0)
        if nid in visited:
            continue
        visited.add(nid)
        levels.setdefault(level, []).append(nid)
        for child in children_map.get(nid, []):
            queue.append((child, level + 1))

    # Add unvisited nodes
    for nid in node_ids:
        if nid not in visited:
            max_level = max(levels.keys()) + 1 if levels else 0
            levels.setdefault(max_level, []).append(nid)

    # Position
    num_levels = len(levels)
    node_w = 2.0
    node_h = 0.7
    v_gap = (_WORK_H - num_levels * node_h) / max(num_levels + 1, 1)

    positions = {}
    for level, nids in levels.items():
        h_gap = (_WORK_W - len(nids) * node_w) / max(len(nids) + 1, 1)
        for j, nid in enumerate(nids):
            x = _LEFT_MARGIN + h_gap * (j + 1) + node_w * j
            y = _TOP_MARGIN + v_gap * (level + 1) + node_h * level
            positions[nid] = (x, y)

    return positions, node_w, node_h


def _layout_block_diagram(nodes: list, edges: list):
    """Auto-layout for block diagram: left-to-right flow."""
    n = len(nodes)
    if n == 0:
        return {}, 0, 0

    rows = min(3, max(1, math.ceil(math.sqrt(n))))
    cols = math.ceil(n / rows)

    node_w = min(2.0, _WORK_W / cols - 0.4)
    node_h = min(1.2, _WORK_H / rows - 0.4)
    h_gap = (_WORK_W - cols * node_w) / max(cols + 1, 1)
    v_gap = (_WORK_H - rows * node_h) / max(rows + 1, 1)

    positions = {}
    for i, node in enumerate(nodes):
        col = i % cols
        row = i // cols
        x = _LEFT_MARGIN + h_gap * (col + 1) + node_w * col
        y = _TOP_MARGIN + v_gap * (row + 1) + node_h * row
        positions[node.get("id", str(i))] = (x, y)

    return positions, node_w, node_h


def _layout_cycle(nodes: list):
    """Auto-layout for cycle: arrange in a circle."""
    n = len(nodes)
    if n == 0:
        return {}, 0, 0

    center_x = _LEFT_MARGIN + _WORK_W / 2
    center_y = _TOP_MARGIN + _WORK_H / 2
    radius = min(_WORK_W, _WORK_H) / 2 - 0.8

    node_w = 1.8
    node_h = 0.7

    positions = {}
    for i, node in enumerate(nodes):
        angle = (2 * math.pi * i / n) - math.pi / 2  # Start from top
        x = center_x + radius * math.cos(angle) - node_w / 2
        y = center_y + radius * math.sin(angle) - node_h / 2
        positions[node.get("id", str(i))] = (x, y)

    return positions, node_w, node_h


def render_diagram_on_slide(slide, diagram_data: Dict[str, Any]) -> None:
    """
    Render a diagram directly on a PowerPoint slide using AutoShapes.

    Args:
        slide: The pptx slide object.
        diagram_data: Dict with keys:
            type: "flowchart" | "hierarchy" | "block_diagram" | "cycle"
            nodes: [{"id": str, "label": str, "shape": str}, ...]
            edges: [{"from": str, "to": str, "label": str}, ...]
    """
    diagram_type = diagram_data.get("type", "flowchart")
    nodes = diagram_data.get("nodes", [])
    edges = diagram_data.get("edges", [])

    if not nodes:
        logger.warning("Empty diagram nodes — skipping render.")
        return

    # Layout based on type
    if diagram_type == "hierarchy" or diagram_type == "pyramid":
        positions, node_w, node_h = _layout_hierarchy(nodes, edges)
    elif diagram_type == "block_diagram" or diagram_type == "swimlane" or diagram_type == "matrix":
        positions, node_w, node_h = _layout_block_diagram(nodes, edges)
    elif diagram_type == "cycle" or diagram_type == "radial":
        positions, node_w, node_h = _layout_cycle(nodes)
    else:  # flowchart / default
        positions, node_w, node_h = _layout_flowchart(nodes, edges)

    # Build node id → node map
    node_map = {}
    for node in nodes:
        node_map[node.get("id", "")] = node

    # Draw edges first (behind shapes)
    for edge in edges:
        from_id = edge.get("from", "")
        to_id = edge.get("to", "")
        if from_id in positions and to_id in positions:
            fx, fy = positions[from_id]
            tx, ty = positions[to_id]

            # Connect from bottom-center of source to top-center of target
            x1 = fx + node_w / 2
            y1 = fy + node_h
            x2 = tx + node_w / 2
            y2 = ty

            # For left-to-right layouts, connect right→left
            if abs(fx - tx) > abs(fy - ty):
                if fx < tx:
                    x1 = fx + node_w
                    y1 = fy + node_h / 2
                    x2 = tx
                    y2 = ty + node_h / 2
                else:
                    x1 = fx
                    y1 = fy + node_h / 2
                    x2 = tx + node_w
                    y2 = ty + node_h / 2

            _add_connector(slide, x1, y1, x2, y2, edge.get("label", ""))

    # Draw nodes
    for i, node in enumerate(nodes):
        nid = node.get("id", str(i))
        if nid not in positions:
            continue

        x, y = positions[nid]
        label = node.get("label", nid)
        shape_name = node.get("shape", "rounded_rect")
        shape_type = _shape_type_for(shape_name)
        color = _NODE_COLORS[i % len(_NODE_COLORS)]

        _add_shape_with_text(
            slide, shape_type,
            x, y, node_w, node_h,
            label, fill_color=color,
            font_size=9 if len(label) > 25 else 10,
            bold=True,
        )

    logger.info(
        "Rendered %s diagram: %d nodes, %d edges",
        diagram_type, len(nodes), len(edges),
    )
