"""
AGRA — Professional PPT Generator (v2)
Builds branded .pptx files with 7 slide layouts, ICG logo,
native shape diagrams, matplotlib charts, and extracted document images.

Slide Layouts:
  1. title          — Title slide with logo, subtitle, branding
  2. section_header  — Section transition slide
  3. bullets         — Standard content with bullet points
  4. two_column      — Side-by-side comparison layout
  5. table           — Native pptx data table
  6. diagram         — Auto-rendered diagram (flowchart/hierarchy/block/cycle)
  7. chart           — Matplotlib chart embedded as image
  8. image           — Full slide image (from document or generated)
  9. thank_you       — Closing slide with branding
"""

import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("agra.ppt_gen")

# ── ICG Brand Palette ──
_NAVY = RGBColor(0x0B, 0x10, 0x20)
_DARK_BLUE = RGBColor(0x1E, 0x3A, 0x8A)
_MID_BLUE = RGBColor(0x2C, 0x5F, 0xC9)
_GOLD = RGBColor(0xD4, 0xA5, 0x37)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
_MED_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
_ACCENT_TEAL = RGBColor(0x06, 0xB6, 0xD4)

# ── Paths ──
_ASSETS_DIR = Path(__file__).resolve().parent.parent.parent / "assets"
_LOGO_PATH = _ASSETS_DIR / "icg_logo.png"

# ── Dimensions ──
_SLIDE_W = 10.0   # inches (widescreen)
_SLIDE_H = 7.5


def _set_slide_bg(slide, color: RGBColor = _NAVY):
    """Set solid background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_logo(slide, right_margin: float = 0.3, top_margin: float = 0.15, size: float = 0.55):
    """Place ICG logo in the top-right corner."""
    if not _LOGO_PATH.exists():
        return
    try:
        left = Inches(_SLIDE_W - right_margin - size)
        top = Inches(top_margin)
        slide.shapes.add_picture(
            str(_LOGO_PATH), left, top,
            width=Inches(size), height=Inches(size),
        )
    except Exception as e:
        logger.debug("Could not add logo: %s", e)


def _add_title_bar(slide, title: str, include_logo: bool = True):
    """Add the standard title bar (dark blue rectangle + gold accent + logo)."""
    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(_SLIDE_W), Inches(1.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _DARK_BLUE
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(8.2), Inches(0.8),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = _WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Gold accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.12), Inches(2.0), Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _GOLD
    accent.line.fill.background()

    if include_logo:
        _add_logo(slide)


def _add_footer(slide, text: str = "AI-Generated Draft | Indian Coast Guard | AGRA System | Confidential"):
    """Add footer bar at the bottom."""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.4),
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = _MED_GRAY
    p.alignment = PP_ALIGN.CENTER


def _add_text_box(
    slide, left, top, width, height, text,
    font_size=16, color=_LIGHT_GRAY, bold=False,
    alignment=PP_ALIGN.LEFT,
):
    """Add a text box with formatting."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def _build_title_slide(prs, title: str, subtitle: str = ""):
    """Layout 1: Title slide with centered branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, _DARK_BLUE)

    # Gradient-like effect: darker band at top
    top_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(_SLIDE_W), Inches(2.5),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x44)
    top_band.line.fill.background()

    # ICG logo centered
    if _LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(
                str(_LOGO_PATH),
                Inches(4.25), Inches(0.6),
                width=Inches(1.5), height=Inches(1.5),
            )
        except Exception:
            pass

    # Title
    _add_text_box(
        slide, 1.0, 2.8, 8.0, 1.2, title,
        font_size=34, color=_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        _add_text_box(
            slide, 1.0, 4.2, 8.0, 0.8, subtitle,
            font_size=18, color=_GOLD, alignment=PP_ALIGN.CENTER,
        )

    # Gold divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(5.2), Inches(3.0), Inches(0.03),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _GOLD
    divider.line.fill.background()

    # Footer
    _add_text_box(
        slide, 1.0, 5.8, 8.0, 0.8,
        "Indian Coast Guard Headquarters, New Delhi\nAGRA — AI-Powered Knowledge Management System",
        font_size=10, color=_MED_GRAY, alignment=PP_ALIGN.CENTER,
    )


def _build_section_header(prs, title: str, subtitle: str = ""):
    """Layout 2: Section transition slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)

    # Large centered title
    _add_text_box(
        slide, 1.0, 2.5, 8.0, 1.5, title,
        font_size=36, color=_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Gold accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(4.2), Inches(3.0), Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _GOLD
    accent.line.fill.background()

    if subtitle:
        _add_text_box(
            slide, 1.0, 4.5, 8.0, 0.8, subtitle,
            font_size=16, color=_LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
        )

    _add_logo(slide)
    _add_footer(slide)


def _build_bullets_slide(prs, title: str, bullets: List[str], notes: str = ""):
    """Layout 3: Standard content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    if bullets:
        txBox = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5), Inches(8.5), Inches(5.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"●  {bullet}"
            p.font.size = Pt(15)
            p.font.color.rgb = _LIGHT_GRAY
            p.space_after = Pt(10)

    _add_footer(slide)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_two_column_slide(prs, title: str, left_data: dict, right_data: dict, notes: str = ""):
    """Layout 4: Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    # Left column header
    left_title = left_data.get("header", "Left")
    left_items = left_data.get("items", [])

    # Left header box
    lh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.5),
    )
    lh.fill.solid()
    lh.fill.fore_color.rgb = _DARK_BLUE
    lh.line.fill.background()
    lh.text_frame.paragraphs[0].text = left_title
    lh.text_frame.paragraphs[0].font.size = Pt(14)
    lh.text_frame.paragraphs[0].font.color.rgb = _GOLD
    lh.text_frame.paragraphs[0].font.bold = True
    lh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left items
    ltb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(4.0), Inches(4.5))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = _LIGHT_GRAY
        p.space_after = Pt(8)

    # Vertical divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.95), Inches(1.5), Inches(0.02), Inches(5.0),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = _GOLD
    div.line.fill.background()

    # Right column
    right_title = right_data.get("header", "Right")
    right_items = right_data.get("items", [])

    rh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.5), Inches(4.2), Inches(0.5),
    )
    rh.fill.solid()
    rh.fill.fore_color.rgb = _DARK_BLUE
    rh.line.fill.background()
    rh.text_frame.paragraphs[0].text = right_title
    rh.text_frame.paragraphs[0].font.size = Pt(14)
    rh.text_frame.paragraphs[0].font.color.rgb = _GOLD
    rh.text_frame.paragraphs[0].font.bold = True
    rh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    rtb = slide.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(4.0), Inches(4.5))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = _LIGHT_GRAY
        p.space_after = Pt(8)

    _add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_table_slide(prs, title: str, table_data: dict, notes: str = ""):
    """Layout 5: Native pptx table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    headers = table_data.get("headers", [])
    rows_data = table_data.get("rows", [])

    if not headers or not rows_data:
        _add_text_box(slide, 1.0, 3.0, 8.0, 1.0, "[No table data provided]",
                       font_size=16, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)
        _add_footer(slide)
        return

    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers)

    # Calculate table dimensions
    table_width = min(9.0, num_cols * 2.0)
    table_left = (_SLIDE_W - table_width) / 2
    row_height = min(0.5, 4.5 / num_rows)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(table_left), Inches(1.5),
        Inches(table_width), Inches(row_height * num_rows),
    )
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = _GOLD
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    # Style data rows
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            if j >= num_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            # Alternating row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2E)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0E, 0x14, 0x24)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = _LIGHT_GRAY
                p.alignment = PP_ALIGN.LEFT

    _add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_diagram_slide(prs, title: str, diagram_data: dict, notes: str = ""):
    """Layout 6: Diagram slide rendered with AutoShapes."""
    from api.generators.diagram_renderer import render_diagram_on_slide

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    try:
        render_diagram_on_slide(slide, diagram_data)
    except Exception as e:
        logger.error("Diagram render failed: %s", e, exc_info=True)
        _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                       f"[Diagram rendering error: {e}]",
                       font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_chart_slide(prs, title: str, chart_data: dict, notes: str = ""):
    """Layout 7: Chart slide with matplotlib-rendered image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    try:
        from api.generators.chart_renderer import render_chart
        chart_path = render_chart(chart_data)

        if chart_path and Path(chart_path).exists():
            # Center the chart image on the slide
            slide.shapes.add_picture(
                chart_path,
                Inches(0.8), Inches(1.5),
                width=Inches(8.4), height=Inches(4.8),
            )
        else:
            _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                           "[Chart could not be rendered]",
                           font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)
    except ImportError:
        logger.warning("matplotlib not available — chart skipped.")
        _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                       "[matplotlib not installed — chart unavailable]",
                       font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)
    except Exception as e:
        logger.error("Chart slide failed: %s", e, exc_info=True)
        _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                       f"[Chart error: {e}]",
                       font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_image_slide(prs, title: str, image_path: str, caption: str = "", notes: str = ""):
    """Layout 8: Full image slide (from uploaded document or generated)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    if image_path and Path(image_path).exists():
        try:
            # Calculate image placement (centered, preserving aspect ratio)
            from PIL import Image
            with Image.open(image_path) as img:
                img_w, img_h = img.size

            aspect = img_w / img_h
            max_w = 8.0
            max_h = 4.5

            if aspect > max_w / max_h:
                disp_w = max_w
                disp_h = max_w / aspect
            else:
                disp_h = max_h
                disp_w = max_h * aspect

            left = (_SLIDE_W - disp_w) / 2
            top = 1.5 + (max_h - disp_h) / 2

            slide.shapes.add_picture(
                image_path,
                Inches(left), Inches(top),
                width=Inches(disp_w), height=Inches(disp_h),
            )
        except Exception as e:
            logger.error("Image insertion failed: %s", e)
            _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                           f"[Image load error: {e}]",
                           font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)
    else:
        _add_text_box(slide, 1.0, 3.0, 8.0, 1.0,
                       "[Image not found]",
                       font_size=14, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Caption
    if caption:
        _add_text_box(slide, 1.0, 6.2, 8.0, 0.5, caption,
                       font_size=11, color=_MED_GRAY, alignment=PP_ALIGN.CENTER)

    _add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _build_sources_slide(prs, sources: List[str], title: str = "Sources & References"):
    """Layout: Sources slide listing referenced documents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _NAVY)
    _add_title_bar(slide, title)

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.5), Inches(5.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, src in enumerate(sources[:12]):  # Cap at 12 to avoid overflow
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"[{i + 1}]  {src}"
        p.font.size = Pt(13)
        p.font.color.rgb = _LIGHT_GRAY
        p.space_after = Pt(10)

    if len(sources) > 12:
        p = tf.add_paragraph()
        p.text = f"... and {len(sources) - 12} more sources."
        p.font.size = Pt(11)
        p.font.color.rgb = _MED_GRAY
        p.space_after = Pt(10)

    _add_footer(slide)


def _build_thank_you_slide(prs, title: str = "Thank You", subtitle: str = ""):
    """Layout 9: Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _DARK_BLUE)

    # Logo centered
    if _LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(
                str(_LOGO_PATH),
                Inches(4.0), Inches(1.2),
                width=Inches(2.0), height=Inches(2.0),
            )
        except Exception:
            pass

    _add_text_box(
        slide, 1.0, 3.5, 8.0, 1.0, title,
        font_size=36, color=_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
    )

    sub = subtitle or "Indian Coast Guard Headquarters\nवयम् रक्षामः — We Protect"
    _add_text_box(
        slide, 1.0, 4.8, 8.0, 1.0, sub,
        font_size=14, color=_GOLD, alignment=PP_ALIGN.CENTER,
    )

    # Gold divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(6.0), Inches(3.0), Inches(0.03),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _GOLD
    divider.line.fill.background()

    _add_text_box(
        slide, 1.0, 6.3, 8.0, 0.5,
        "Generated by AGRA — AI-Powered Knowledge Management System",
        font_size=9, color=_MED_GRAY, alignment=PP_ALIGN.CENTER,
    )


# ═══════════════════════════════════════════════════════════════
#  MAIN BUILD FUNCTION
# ═══════════════════════════════════════════════════════════════

def build_pptx(
    slides_data: List[Dict[str, Any]],
    output_path: str,
    title: str = "AGRA Presentation",
    template_path: Optional[str] = None,
    extracted_images: Optional[List[Dict[str, Any]]] = None,
) -> str:
    """
    Build a professional .pptx file from structured slide data.

    Args:
        slides_data: List of slide dicts. Each must have:
            - "title": str
            - "layout": str — one of: "title", "section_header", "bullets",
              "two_column", "table", "diagram", "chart", "image", "thank_you"
            - Layout-specific fields (see individual builders)
        output_path: Where to save the .pptx file.
        title: Presentation title.
        template_path: Optional path to a .pptx template.
        extracted_images: Optional list of extracted document images.

    Returns:
        Path to the generated .pptx file.
    """
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(_SLIDE_W)
        prs.slide_height = Inches(_SLIDE_H)

    image_index = 0  # Track extracted images for insertion

    for i, sd in enumerate(slides_data):
        layout = sd.get("layout", "bullets")
        slide_title = sd.get("title", f"Slide {i + 1}")
        bullets = sd.get("bullets", [])
        notes = sd.get("notes", "")

        try:
            if layout == "title" or i == 0:
                subtitle = sd.get("subtitle", bullets[0] if bullets else "")
                _build_title_slide(prs, slide_title, subtitle)

            elif layout == "section_header":
                subtitle = sd.get("subtitle", bullets[0] if bullets else "")
                _build_section_header(prs, slide_title, subtitle)

            elif layout == "two_column":
                left_data = sd.get("left_column", {"header": "Left", "items": []})
                right_data = sd.get("right_column", {"header": "Right", "items": []})
                _build_two_column_slide(prs, slide_title, left_data, right_data, notes)

            elif layout == "table":
                table_data = sd.get("table_data", {})
                _build_table_slide(prs, slide_title, table_data, notes)

            elif layout == "diagram":
                diagram_data = sd.get("diagram_data", {})
                _build_diagram_slide(prs, slide_title, diagram_data, notes)

            elif layout == "chart":
                chart_data = sd.get("chart_data", {})
                _build_chart_slide(prs, slide_title, chart_data, notes)

            elif layout == "image":
                img_path = sd.get("image_path", "")
                # Auto-fill from extracted document images
                if not img_path and extracted_images and image_index < len(extracted_images):
                    img_path = extracted_images[image_index].get("path", "")
                    image_index += 1
                caption = sd.get("caption", "")
                _build_image_slide(prs, slide_title, img_path, caption, notes)

            elif layout == "sources":
                sources = sd.get("sources", [])
                _build_sources_slide(prs, sources, slide_title)

            elif layout == "thank_you":
                subtitle = sd.get("subtitle", "")
                _build_thank_you_slide(prs, slide_title, subtitle)

            else:
                # Default: bullets layout
                _build_bullets_slide(prs, slide_title, bullets, notes)

        except Exception as e:
            logger.error("Failed to build slide %d (%s): %s", i, layout, e, exc_info=True)
            # Fallback: render as bullets
            try:
                _build_bullets_slide(prs, slide_title, bullets, notes)
            except Exception:
                pass

    prs.save(output_path)
    logger.info("PPTX saved: %s (%d slides)", output_path, len(slides_data))
    return output_path
